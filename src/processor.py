import fitz  # PyMuPDF
import os
from PIL import Image, ImageFilter
import io
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from .config import Config

from rapidocr_onnxruntime import RapidOCR
import logging
import numpy as np

# Suppress PaddleOCR logging
logging.getLogger("ppocr").setLevel(logging.ERROR)

class PDFProcessor:
    def __init__(self, input_path, output_dir=None, font_path=None):
        self.input_path = input_path
        self.output_dir = output_dir or Config.OUTPUT_DIR
        
        # Ensure output directory exists
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir, exist_ok=True)

        self.font_path = font_path or Config.DEFAULT_FONT_PATH
        
        if not os.path.exists(self.font_path):
            if os.path.exists(Config.SYSTEM_FONT_FALLBACK):
                 self.font_path = Config.SYSTEM_FONT_FALLBACK
            else:
                print(f"Warning: Font file not found at {self.font_path}. Text rendering might fail or use default.")

        self.doc = fitz.open(self.input_path)
        self.filename = os.path.splitext(os.path.basename(input_path))[0]
        self.ocr = None # Lazy init

    def _init_ocr(self):
        if self.ocr is None:
            # Initialize PaddleOCR
            # Disable angle classifier to avoid "unexpected keyword argument 'cls'" error
            # NotebookLM slides are usually horizontal anyway
            print("Initializing RapidOCR...")
            self.ocr = RapidOCR()

    def get_page_thumbnails(self, dpi=72):
        """
        Generates thumbnails for all pages.
        Returns a list of tuples: (page_num, pil_image)
        """
        thumbnails = []
        for page_num in range(len(self.doc)):
            page = self.doc[page_num]
            pix = page.get_pixmap(dpi=dpi)
            img = Image.open(io.BytesIO(pix.tobytes()))
            thumbnails.append((page_num + 1, img))
        return thumbnails

    def _apply_watermark_removal(self, img, wm_settings):
        """
        Applies watermark removal to the given image based on settings.
        """
        if not wm_settings:
            return img
            
        w, h = img.size
        x_start = int(w * wm_settings["x_start"])
        y_start = int(h * wm_settings["y_start"])
        width = int(w * wm_settings["width"])
        height = int(h * wm_settings["height"])
        
        if wm_settings.get("use_mirror_patch", False):
            # Mirror Patch mode
            src_x = w - (x_start + width)
            src_y = y_start
            src_x = max(0, int(src_x))
            
            patch = img.crop((src_x, src_y, src_x + width, src_y + height))
            patch = patch.transpose(Image.FLIP_LEFT_RIGHT)
            img.paste(patch, (x_start, y_start))
            
        elif wm_settings.get("use_patch", False):
            # Manual Patch mode
            src_x = int(w * wm_settings["src_x"])
            src_y = int(h * wm_settings["src_y"])
            
            src_x = max(0, min(src_x, w - width))
            src_y = max(0, min(src_y, h - height))
            
            patch = img.crop((src_x, src_y, src_x + width, src_y + height))
            img.paste(patch, (x_start, y_start))
        else:
            # White mask mode
            img.paste((255, 255, 255), (x_start, y_start, x_start + width, y_start + height))
            
        return img

    def extract_elements(self, page_num, enable_ocr=False):
        """
        Extracts text blocks from a page.
        Returns a list of dictionaries containing text, bbox, size, color.
        """
        page = self.doc[page_num]
        text_instances = []
        
        # 1. Try standard PDF extraction first
        blocks = page.get_text("dict")["blocks"]
        
        for block in blocks:
            if "lines" in block:
                for line in block["lines"]:
                    for span in line["spans"]:
                        text = span["text"].strip()
                        if not text:
                            continue
                        
                        # Convert sRGB int to Hex
                        srgb = span["color"]
                        r = (srgb >> 16) & 255
                        g = (srgb >> 8) & 255
                        b = srgb & 255
                        hex_color = "#{:02x}{:02x}{:02x}".format(r, g, b)
                        
                        text_instances.append({
                            "text": text,
                            "bbox": span["bbox"], # (x0, y0, x1, y1)
                            "size": span["size"],
                            "color": hex_color, # Hex string
                            "origin": span["origin"]
                        })
        
        # 2. OCR Fallback
        # If no text found, or very little text (e.g. just page numbers), try OCR
        if enable_ocr and len(text_instances) < 5: # Threshold can be adjusted
            print(f"Page {page_num}: Low text count ({len(text_instances)}). Attempting OCR...")
            self._init_ocr()
            
            # Get page image for OCR
            # Use 150 DPI (down from 200) to improve speed while maintaining acceptable accuracy
            pix = page.get_pixmap(dpi=150)
            
            # Convert to numpy array for RapidOCR
            # pix.samples is bytes
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.height, pix.width, pix.n)
            
            # Handle Alpha channel if present
            if pix.n == 4:
                img = np.ascontiguousarray(img[..., :3]) # Drop alpha, keep RGB
            elif pix.n == 3:
                pass # Already RGB
            
            # Run OCR with RapidOCR
            result, elapse = self.ocr(img)
            
            if result:
                scale_x = page.rect.width / pix.width
                scale_y = page.rect.height / pix.height
                
                for item in result:
                    # item structure: [dt_box, text, score]
                    # dt_box: [[x1, y1], [x2, y2], [x3, y3], [x4, y4]]
                    dt_box, text, score = item
                    
                    if score < 0.5: continue
                    
                    # Calculate bbox from dt_box
                    xs = [p[0] for p in dt_box]
                    ys = [p[1] for p in dt_box]
                    x0 = min(xs) * scale_x
                    y0 = min(ys) * scale_y
                    x1 = max(xs) * scale_x
                    y1 = max(ys) * scale_y
                    
                    # Estimate font size
                    size = (y1 - y0) * 0.8
                    
                    # Sample color from image
                    # Convert bbox to pixel coordinates for sampling
                    ix0 = int(x0 / scale_x)
                    iy0 = int(y0 / scale_y)
                    ix1 = int(x1 / scale_x)
                    iy1 = int(y1 / scale_y)
                    
                    # Ensure bounds
                    h_img, w_img, _ = img.shape
                    ix0 = max(0, ix0)
                    iy0 = max(0, iy0)
                    ix1 = min(w_img, ix1)
                    iy1 = min(h_img, iy1)
                    
                    if ix1 > ix0 and iy1 > iy0:
                        # Crop region
                        region = img[iy0:iy1, ix0:ix1]
                        # Get average color (simple mean)
                        avg_color = region.mean(axis=(0, 1)).astype(int) # RGB
                        hex_color = "#{:02x}{:02x}{:02x}".format(avg_color[0], avg_color[1], avg_color[2])
                    else:
                        hex_color = "#000000" # Fallback to black

                    text_instances.append({
                        "text": text,
                        "bbox": (x0, y0, x1, y1),
                        "size": size,
                        "color": hex_color,
                        "origin": (x0, y1)
                    })
        
        return text_instances
    def clean_page_image(self, page_num, dpi=300, wm_settings=None):
        """
        Renders the page as an image and removes the watermark.
        Returns a PIL Image object.
        """
        page = self.doc[page_num]
        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes()))
        
        if wm_settings:
            img = self._apply_watermark_removal(img, wm_settings)
        
        return img

    def get_background_image(self, page_num, dpi=300, wm_settings=None):
        """
        Hides text, renders page to image (background only), then restores text.
        """
        # Open a fresh handle to avoid messing up the main doc state if we were to modify it
        doc_bg = fitz.open(self.input_path)
        page_bg = doc_bg[page_num]
        
        # Redact all text
        text_blocks = page_bg.get_text("blocks")
        for block in text_blocks:
            # block: (x0, y0, x1, y1, "text", block_no, block_type)
            if block[6] == 0: # Text block
                rect = fitz.Rect(block[:4])
                # Add redaction annotation
                # fill=None means no fill color (transparent/white depending on viewer, but usually removes content)
                page_bg.add_redact_annot(rect)
                
        # Apply redactions
        # images=fitz.PDF_REDACT_IMAGE_NONE ensures we DON'T remove images that might be under the text
        # graphics=fitz.PDF_REDACT_IMAGE_NONE ensures we keep vector graphics (lines etc)
        page_bg.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE, graphics=fitz.PDF_REDACT_IMAGE_NONE)
        
        # Now render
        pix = page_bg.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes()))
        
        doc_bg.close()
        
        # Apply Watermark Mask
        # Apply Watermark Mask
        if wm_settings:
            img = self._apply_watermark_removal(img, wm_settings)
        
        return img

    def render_new_pdf(self, wm_settings=None, debug_mode=False, enable_ocr=False, progress_callback=None, pages_to_remove=None):
        """
        Creates a new PDF with high-quality text and original background.
        pages_to_remove: List of 0-based page numbers to skip.
        """
        output_path = os.path.join(self.output_dir, f"{self.filename}_enhanced.pdf")
        new_doc = fitz.open()
        
        # Verify font availability
        try:
            test_page = new_doc.new_page()
            test_page.insert_font(fontname="test_font", fontfile=self.font_path)
        except Exception as e:
            print(f"Error loading font {self.font_path}: {e}")
        finally:
            # Ensure the test page is removed regardless of success or failure
            if len(new_doc) > 0:
                new_doc.delete_page(0)
        
        total_pages = len(self.doc)
        for page_num in range(total_pages):
            if pages_to_remove and page_num in pages_to_remove:
                continue

            if progress_callback:
                progress_callback(page_num / total_pages, f"Processing page {page_num + 1}/{total_pages}")

            # 1. Get Text
            text_elements = self.extract_elements(page_num, enable_ocr=enable_ocr)
            
            # 2. Get Background (Cleaned)
            # Use JPEG to reduce size significantly
            bg_img = self.get_background_image(page_num, dpi=Config.DPI, wm_settings=wm_settings)
            
            # Save bg to temp file to insert into PDF
            # Use JPEG with quality 80
            temp_bg_path = os.path.join(self.output_dir, f"temp_bg_{page_num}.jpg")
            bg_img.save(temp_bg_path, "JPEG", quality=80, optimize=True)
            
            # 3. Create New Page
            page = self.doc[page_num]
            new_page = new_doc.new_page(width=page.rect.width, height=page.rect.height)
            
            # 4. Insert Background
            new_page.insert_image(new_page.rect, filename=temp_bg_path)
            
            # 5. Insert Text
            fontname = "custom_font"
            try:
                new_page.insert_font(fontname=fontname, fontfile=self.font_path)
            except Exception:
                # If font fails, it might use default.
                pass
            
            for elem in text_elements:
                hex_color = elem["color"]
                if hex_color.startswith("#"):
                    hex_color = hex_color[1:]
                
                try:
                    r = int(hex_color[0:2], 16) / 255.0
                    g = int(hex_color[2:4], 16) / 255.0
                    b = int(hex_color[4:6], 16) / 255.0
                except ValueError:
                    r, g, b = 0, 0, 0
                
                # Debug Mode: Force Red Color
                if debug_mode:
                    r, g, b = (1, 0, 0)
                
                # Draw Text Background (to cover old blurry text)
                # We use a simple white box for now. 
                # Ideally we could pick the average color of the background in that rect.
                if wm_settings and wm_settings.get("text_bg", False):
                    # Draw a filled rectangle
                    # bbox is (x0, y0, x1, y1)
                    rect = fitz.Rect(elem["bbox"])
                    # Add a small padding
                    rect.x0 -= 1
                    rect.y0 -= 1
                    rect.x1 += 1
                    rect.y1 += 1
                    
                    # Draw white box
                    shape = new_page.new_shape()
                    shape.draw_rect(rect)
                    shape.finish(color=None, fill=(1, 1, 1)) # White fill, no border
                    shape.commit()

                new_page.insert_text(
                    elem["origin"],
                    elem["text"],
                    fontname=fontname,
                    fontsize=elem["size"],
                    color=(r, g, b)
                )
            
            # Cleanup temp file
            if os.path.exists(temp_bg_path):
                os.remove(temp_bg_path)
        
        if progress_callback:
            progress_callback(1.0, "PDF generation complete!")

        new_doc.save(output_path)
        print(f"PDF saved to: {output_path}")
        return output_path

    def convert_to_pptx(self, wm_settings=None, text_mode="re-render", enable_ocr=False, progress_callback=None, pages_to_remove=None):
        """
        Converts the PDF to a PPTX file with editable text.
        text_mode: 're-render' (clean bg + new text) or 'overlay' (original bg + invisible text)
        pages_to_remove: List of 0-based page numbers to skip.
        """
        output_path = os.path.join(self.output_dir, f"{self.filename}.pptx")
        prs = Presentation()
        
        first_page = self.doc[0]
        # PPTX uses EMU (English Metric Unit). 1 point = 12700 EMUs.
        # Ensure we don't overflow or create tiny slides.
        prs.slide_width = int(first_page.rect.width * 12700)
        prs.slide_height = int(first_page.rect.height * 12700)
        
        total_pages = len(self.doc)
        for page_num in range(total_pages):
            if pages_to_remove and page_num in pages_to_remove:
                continue

            if progress_callback:
                progress_callback(page_num / total_pages, f"Converting page {page_num + 1}/{total_pages}")

            # 1. Get Background
            # Use lower DPI for PPTX background to keep file light
            if text_mode == "overlay":
                # Overlay mode: Use original image (cleaned of watermark only)
                bg_img = self.clean_page_image(page_num, dpi=150, wm_settings=wm_settings)
            else:
                # Re-render mode: Use redacted background (text removed)
                bg_img = self.get_background_image(page_num, dpi=150, wm_settings=wm_settings) 
                
            temp_bg_path = os.path.join(self.output_dir, f"temp_bg_pptx_{page_num}.jpg")
            bg_img.save(temp_bg_path, "JPEG", quality=80)
            
            # 2. Add Slide
            blank_slide_layout = prs.slide_layouts[6] 
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # 3. Set Background
            left = top = 0
            pic = slide.shapes.add_picture(temp_bg_path, left, top, width=prs.slide_width, height=prs.slide_height)
            
            # 4. Add Text Boxes
            text_elements = self.extract_elements(page_num, enable_ocr=enable_ocr)
            
            # Normalize font sizes to standard PPTX sizes
            text_elements = self._normalize_font_sizes(text_elements)
            
            if not text_elements:
                print(f"Warning: No text found on page {page_num}. PPTX slide will be image only.")
            
            for elem in text_elements:
                x, y, x1, y1 = elem["bbox"]
                w = x1 - x
                h = y1 - y
                
                # Minimum size check to avoid tiny boxes
                if w < 1 or h < 1:
                    continue

                # Add text box
                txBox = slide.shapes.add_textbox(Pt(x), Pt(y), Pt(w), Pt(h))
                tf = txBox.text_frame
                tf.word_wrap = True 
                
                p = tf.paragraphs[0]
                p.text = "" # Clear paragraph text
                run = p.add_run()
                run.text = elem["text"]
                run.font.size = Pt(elem["size"])
                run.font.name = "Microsoft JhengHei"
                
                if text_mode == "overlay":
                    # Overlay mode: Invisible text
                    # Set color to black but fully transparent
                    # NOTE: We use background() to set "No Fill", making text invisible.
                    run.font.fill.background() 
                else:
                    # Re-render mode: Visible colored text
                    hex_color = elem["color"]
                    if hex_color.startswith("#"):
                        hex_color = hex_color[1:]
                    
                    try:
                        r = int(hex_color[0:2], 16)
                        g = int(hex_color[2:4], 16)
                        b = int(hex_color[4:6], 16)
                    except ValueError:
                        r, g, b = 0, 0, 0 # Fallback to black

                    run.font.color.rgb = RGBColor(r, g, b)

            # 5. Add Notes (Speaker Notes)
            if text_elements:
                # Combine all text into a single string
                notes_text = "\n".join([elem["text"] for elem in text_elements])
                
                # Access notes slide (creates it if not exists)
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = notes_text

            # Cleanup
            if os.path.exists(temp_bg_path):
                os.remove(temp_bg_path)
        
        if progress_callback:
            progress_callback(1.0, "PPTX conversion complete!")
            
        prs.save(output_path)
        print(f"PPTX saved to: {output_path}")
        return output_path

    def extract_text_data(self, pages=None, progress_callback=None):
        """
        Extracts text data from the PDF.
        pages: Optional list of 1-based page numbers to process. If None, process all.
        """
        all_text_data = []
        
        # Determine pages to process
        if pages:
            # Convert 1-based to 0-based
            pages_to_process = [p - 1 for p in pages]
        else:
            pages_to_process = range(len(self.doc))
            
        total_pages = len(pages_to_process)
        for idx, page_num in enumerate(pages_to_process):
            if progress_callback:
                progress_callback(idx / total_pages, f"Analyzing page {page_num + 1}")

            if page_num < 0 or page_num >= len(self.doc):
                continue
                
            elements = self.extract_elements(page_num, enable_ocr=True)
            for i, elem in enumerate(elements):
                # Create a unique ID for each text element
                elem_id = f"p{page_num}_e{i}"
                
                all_text_data.append({
                    "id": elem_id,
                    "page": page_num + 1, # 1-based for display
                    "original_text": elem["text"],
                    "new_text": elem["text"], # Default to original
                    "bbox": elem["bbox"],
                    "size": elem["size"],
                    "color": elem["color"],
                    "origin": elem["origin"]
                })
        
        if progress_callback:
            progress_callback(1.0, "Analysis complete!")
                
        # Normalize font sizes across the document
        all_text_data = self._normalize_font_sizes(all_text_data)
        
        # Normalize coordinates (alignment snapping)
        all_text_data = self._normalize_coordinates(all_text_data)
        
        # Heuristic BBox Adjustment (Ensure bbox covers text)
        all_text_data = self._adjust_bbox_by_content(all_text_data)
        
        return all_text_data
        


    def _normalize_font_sizes(self, text_data):
        """
        Normalizes font sizes by clustering similar sizes and snapping to standard Word sizes.
        """
        if not text_data:
            return text_data
            
        # Standard Word font sizes
        STANDARD_SIZES = [8, 9, 10, 10.5, 11, 12, 14, 16, 18, 20, 24, 28, 32, 36, 40, 44, 48, 54, 60, 66, 72, 80, 88, 96]
        
        # Extract all sizes
        sizes = sorted([item["size"] for item in text_data])
        if not sizes:
            return text_data
            
        # Cluster sizes
        clusters = []
        current_cluster = [sizes[0]]
        
        for s in sizes[1:]:
            # If size is within 1.5pt of the cluster start, add to cluster
            # Using tighter tolerance (1.5) to distinguish close sizes like 10, 10.5, 11
            if s - current_cluster[0] < 1.5:
                current_cluster.append(s)
            else:
                clusters.append(current_cluster)
                current_cluster = [s]
        clusters.append(current_cluster)
        
        # Create mapping
        size_map = {}
        for cluster in clusters:
            # Calculate average
            avg = sum(cluster) / len(cluster)
            
            # Find closest standard size
            closest_size = min(STANDARD_SIZES, key=lambda x: abs(x - avg))
            
            # If the difference is reasonable (e.g. within 2.0), snap to it.
            # Otherwise, just round to nearest 0.5
            if abs(closest_size - avg) <= 2.0:
                 norm_size = closest_size
            else:
                 norm_size = round(avg * 2) / 2
            
            for s in cluster:
                size_map[s] = norm_size
                
        # Apply mapping
        for item in text_data:
            item["size"] = size_map.get(item["size"], item["size"])
            
        return text_data

    def _normalize_coordinates(self, text_data):
        """
        Normalizes coordinates (bbox X/Y) by clustering similar values.
        """
        if not text_data:
            return text_data
            
        # Helper to cluster and map values
        def get_snap_map(values, tolerance=2.0):
            if not values: return {}
            sorted_vals = sorted(values)
            clusters = []
            current_cluster = [sorted_vals[0]]
            
            for v in sorted_vals[1:]:
                if v - current_cluster[0] < tolerance:
                    current_cluster.append(v)
                else:
                    clusters.append(current_cluster)
                    current_cluster = [v]
            clusters.append(current_cluster)
            
            snap_map = {}
            for cluster in clusters:
                avg = sum(cluster) / len(cluster)
                # Round to nearest 0.5 for clean numbers
                snapped = round(avg * 2) / 2
                for v in cluster:
                    snap_map[v] = snapped
            return snap_map

        # Collect all coordinates
        x0s = [item["bbox"][0] for item in text_data]
        y0s = [item["bbox"][1] for item in text_data]
        x1s = [item["bbox"][2] for item in text_data]
        y1s = [item["bbox"][3] for item in text_data]
        
        # Create maps
        map_x0 = get_snap_map(x0s, tolerance=3.0) # Slightly larger tolerance for alignment
        map_y0 = get_snap_map(y0s, tolerance=3.0)
        map_x1 = get_snap_map(x1s, tolerance=3.0)
        map_y1 = get_snap_map(y1s, tolerance=3.0)
        
        # Apply snapping
        for item in text_data:
            b = item["bbox"]
            new_x0 = map_x0.get(b[0], b[0])
            new_y0 = map_y0.get(b[1], b[1])
            new_x1 = map_x1.get(b[2], b[2])
            new_y1 = map_y1.get(b[3], b[3])
            
            item["bbox"] = (new_x0, new_y0, new_x1, new_y1)
            
            # Update origin (usually x0, y1 for bottom-left, but PyMuPDF origin varies)
            # We extracted origin from span["origin"]. 
            # Let's assume origin.x should align with bbox.x0 and origin.y with bbox.y1 (baseline)
            # But baseline is not exactly y1. 
            # Let's just snap the origin X to x0 if they were close, and origin Y to y1 if close?
            # Actually, safer to just snap origin X/Y independently if we tracked them.
            # Since we didn't track origin separately for clustering, let's just shift it 
            # by the same delta as the bbox shift.
            
            dx = new_x0 - b[0]
            dy = new_y1 - b[3] # Use bottom shift for Y? Or top? Text is usually baseline aligned.
            
            # Let's try to snap origin.x to new_x0 directly if they are close
            ox, oy = item["origin"]
            new_ox = ox + dx
            new_oy = oy + dy # This might be risky if y1 changed differently than baseline.
            
            item["origin"] = (new_ox, new_oy)
            
        return text_data

    def _adjust_bbox_by_content(self, text_data):
        """
        Adjusts bbox size based on heuristic calculation of text content.
        Ensures the bbox is large enough to cover the text.
        """
        if not text_data:
            return text_data
            
        for item in text_data:
            text = item["original_text"]
            font_size = item["size"]
            bbox = item["bbox"]
            x0, y0, x1, y1 = bbox
            
            current_w = x1 - x0
            current_h = y1 - y0
            
            # Calculate expected width
            # Chinese/Full-width: ~1.0 * font_size
            # ASCII: ~0.6 * font_size
            ascii_count = 0
            non_ascii_count = 0
            for char in text:
                if ord(char) < 128:
                    ascii_count += 1
                else:
                    non_ascii_count += 1
            
            expected_w = (ascii_count * 0.6 + non_ascii_count * 1.0) * font_size
            # Add some safety margin (10%)
            expected_w *= 1.1
            
            # Calculate expected height
            # Line height is usually 1.2 * font_size
            expected_h = font_size * 1.2
            
            # Check and expand
            new_x0, new_y0, new_x1, new_y1 = x0, y0, x1, y1
            
            if current_w < expected_w:
                # Expand horizontally from center
                center_x = (x0 + x1) / 2
                half_w = expected_w / 2
                new_x0 = center_x - half_w
                new_x1 = center_x + half_w
                
            if current_h < expected_h:
                # Expand vertically from center
                center_y = (y0 + y1) / 2
                half_h = expected_h / 2
                new_y0 = center_y - half_h
                new_y1 = center_y + half_h
            item["bbox"] = (new_x0, new_y0, new_x1, new_y1)
            
        return text_data

    def process_background_regions(self, page_num, text_bboxes, dpi=300, padding=5, mode='Blur'):
        """
        Processes background regions for text replacement.
        mode: 'Blur', 'Smart Fill', 'White'
        """
        page = self.doc[page_num]
        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes()))
        
        w, h = img.size
        scale_x = w / page.rect.width
        scale_y = h / page.rect.height
        
        for bbox in text_bboxes:
            # Scale bbox to image coordinates
            x0, y0, x1, y1 = bbox
            ix0 = int(x0 * scale_x)
            iy0 = int(y0 * scale_y)
            ix1 = int(x1 * scale_x)
            iy1 = int(y1 * scale_y)
            
            # Apply padding
            px = int(padding * scale_x)
            py = int(padding * scale_y)
            
            ix0 -= px
            iy0 -= py
            ix1 += px
            iy1 += py
            
            # Ensure bounds
            ix0 = max(0, ix0)
            iy0 = max(0, iy0)
            ix1 = min(w, ix1)
            iy1 = min(h, iy1)
            
            if ix1 <= ix0 or iy1 <= iy0:
                continue
                
            region_box = (ix0, iy0, ix1, iy1)
            
            if mode == 'White':
                img.paste((255, 255, 255), region_box)
                
            elif mode == 'Smart Fill':
                # Robust approach: Get colors from corners
                region = img.crop(region_box)
                corners = [
                    region.getpixel((0, 0)),
                    region.getpixel((region.width-1, 0)),
                    region.getpixel((0, region.height-1)),
                    region.getpixel((region.width-1, region.height-1))
                ]
                
                # Average r, g, b
                avg_r = sum(c[0] for c in corners) // 4
                avg_g = sum(c[1] for c in corners) // 4
                avg_b = sum(c[2] for c in corners) // 4
                
                fill_color = (avg_r, avg_g, avg_b)
                img.paste(fill_color, region_box)
                
            else: # Blur (Default)
                region = img.crop(region_box)
                # Radius 20 for strong blur
                blurred_region = region.filter(ImageFilter.GaussianBlur(radius=20))
                img.paste(blurred_region, (ix0, iy0))
            
        return img

    def apply_text_edits(self, edits_data, font_path=None, wm_settings=None, bg_mode='Blur'):
        """
        Applies text edits to the PDF.
        edits_data: List of dicts (from st.data_editor)
        wm_settings: Watermark settings to apply to background.
        bg_mode: 'Blur', 'Smart Fill', 'White'
        """
        output_path = os.path.join(self.output_dir, f"{self.filename}_edited.pdf")
        
        # Group edits by page
        edits_by_page = {}
        for item in edits_data:
            page_num = item["page"]
            if page_num not in edits_by_page:
                edits_by_page[page_num] = []
            edits_by_page[page_num].append(item)
            
        # Use a temporary PDF to build the new one
        new_doc = fitz.open()
        
        # Font handling
        fontname = "cjk"
        use_font_path = font_path or self.font_path
        
        # Create font object for measuring text width
        try:
            measure_font = fitz.Font(fontfile=use_font_path)
        except:
            measure_font = fitz.Font("helv") # Fallback
        
        # Process each page of the ORIGINAL document
        for page_num in range(len(self.doc)):
            original_page = self.doc[page_num]
            
            # Check if this page is in our edits data
            if (page_num + 1) not in edits_by_page:
                # Page was not selected for editing/analysis. Copy original.
                new_doc.insert_pdf(self.doc, from_page=page_num, to_page=page_num)
                continue
            
            # Identify modified items for this page
            page_edits = edits_by_page.get(page_num + 1, [])
            modified_items = [item for item in page_edits if item["original_text"] != item["new_text"]]
            
            # 1. Prepare Background
            if modified_items:
                # We have modifications, so we need to process background for those specific items
                bboxes_to_blur = []
                
                for item in modified_items:
                    old_bbox = item["bbox"]
                    new_text = item["new_text"]
                    font_size = item["size"]
                    
                    # Calculate new text dimensions
                    try:
                        new_text_width = measure_font.text_length(new_text, fontsize=font_size)
                    except:
                        new_text_width = (len(new_text) * font_size) # Rough fallback
                        
                    new_text_height = font_size * 1.2 # Approx line height
                    
                    # Calculate Union BBox
                    # Old bbox: (x0, y0, x1, y1)
                    x0, y0, x1, y1 = old_bbox
                    old_w = x1 - x0
                    old_h = y1 - y0
                    
                    # Determine max width and height
                    final_w = max(old_w, new_text_width)
                    final_h = max(old_h, new_text_height)
                    
                    # Align: Left-aligned for X, Center-aligned for Y
                    # New X range: start at x0, extend to x0 + final_w
                    new_x0 = x0
                    new_x1 = x0 + final_w
                    
                    # New Y range: center around old center
                    center_y = (y0 + y1) / 2
                    new_y0 = center_y - (final_h / 2)
                    new_y1 = center_y + (final_h / 2)
                    
                    # Add a little padding to the calculation itself? 
                    # process_background_regions already adds padding, so we just pass the "content" bbox.
                    bboxes_to_blur.append((new_x0, new_y0, new_x1, new_y1))
                
                # Use larger padding (e.g. 5 points)
                bg_img = self.process_background_regions(page_num, bboxes_to_blur, dpi=Config.DPI, padding=5, mode=bg_mode)
            else:
                # No edits for this page, just render standard background
                pix = original_page.get_pixmap(dpi=Config.DPI)
                bg_img = Image.open(io.BytesIO(pix.tobytes()))

            # Apply Watermark Removal if requested
            if wm_settings:
                bg_img = self._apply_watermark_removal(bg_img, wm_settings)

            # Save bg to temp
            temp_bg_path = os.path.join(self.output_dir, f"temp_bg_edit_{page_num}.jpg")
            bg_img.save(temp_bg_path, "JPEG", quality=85)
            
            # 2. Create New Page
            new_page = new_doc.new_page(width=original_page.rect.width, height=original_page.rect.height)
            
            # 3. Insert Background
            new_page.insert_image(new_page.rect, filename=temp_bg_path)
            
            # 4. Insert Text (Iterate ALL items to ensure copyability)
            if page_edits:
                # Register font
                try:
                    new_page.insert_font(fontname=fontname, fontfile=use_font_path)
                except Exception:
                    pass
                
                for item in page_edits:
                    # Check if this item is modified
                    is_modified = item in modified_items
                    
                    # Parse Hex color
                    hex_color = item["color"]
                    if hex_color.startswith("#"):
                        hex_color = hex_color[1:]
                    try:
                        r = int(hex_color[0:2], 16) / 255.0
                        g = int(hex_color[2:4], 16) / 255.0
                        b = int(hex_color[4:6], 16) / 255.0
                    except:
                        r, g, b = 0, 0, 0 # Fallback to black
                    
                    if is_modified:
                        # Modified: Render Visible Text
                        render_mode = 0 # Fill
                        text_content = item["new_text"]
                    else:
                        # Unmodified: Render Invisible Text (for copyability)
                        render_mode = 3 # Invisible
                        text_content = item["original_text"] # Use original text
                    
                    # Insert text
                    new_page.insert_text(
                        item["origin"],
                        text_content,
                        fontname=fontname,
                        fontsize=item["size"],
                        color=(r, g, b),
                        render_mode=render_mode
                    )
            
            # Cleanup
            if os.path.exists(temp_bg_path):
                os.remove(temp_bg_path)
        
        new_doc.save(output_path)
        return output_path
